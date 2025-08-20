
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm

def t_i18n(key, lang, i18n):
    try:
        return i18n.get(lang, {}).get(key, i18n.get('en', {}).get(key, key))
    except Exception:
        return key

def _brand_header(slide, brand_primary="#C00000", logo_path=None):
    try:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.18))
        band.fill.solid()
        h = brand_primary.lstrip('#')
        band.fill.fore_color.rgb = RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
        band.line.fill.background()
        if logo_path and os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(8.2), Inches(0.2), height=Inches(0.6))
    except Exception:
        pass

def _load_brand_master_fallback(template_path):
    """If template_path is None, try to read templates.yaml → brand.pptx_master; else fallback to default Presentation()."""
    master_to_use = template_path
    try:
        import yaml
        if master_to_use is None and os.path.exists('templates.yaml'):
            with open('templates.yaml','r', encoding='utf-8') as f:
                tpl = yaml.safe_load(f) or {}
                master_to_use = (tpl.get('brand',{}) or {}).get('pptx_master')
    except Exception:
        master_to_use = template_path
    try:
        if master_to_use and os.path.exists(master_to_use):
            return Presentation(master_to_use)
        return Presentation()
    except Exception:
        return Presentation()

def add_material_flow_narrative_slide(prs, text: str, lang='en', i18n=None, brand_primary="#C00000", logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _brand_header(slide, brand_primary, logo_path)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = t_i18n("material_flow", lang, i18n or {})
    title.text_frame.paragraphs[0].font.size = Pt(26)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.text = text
    tf.paragraphs[0].font.size = Pt(12)

def export_observations_pptx(observations_df, out_path, steps=None, perstep_top2=None, spacing_mode="Effective CT", ct_eff_map=None, vc_summary=None, material_flow_text=None, photos=None, template_path=None, lang='en', i18n=None, brand_primary="#C00000", logo_path=None, finance=None, product_df=None, champion=None, savings=None):
    prs = _load_brand_master_fallback(template_path)
    title = prs.slides.add_slide(prs.slide_layouts[0])
    _brand_header(title, brand_primary, logo_path)
    if title.shapes.title:
        title.shapes.title.text = t_i18n("title", lang, i18n or {})
    if hasattr(title, "placeholders") and len(title.placeholders)>1:
        try:
            title.placeholders[1].text = ""
        except Exception:
            pass

    if steps and perstep_top2:
        add_current_state_map_slide(prs, steps, perstep_top2, spacing_mode=spacing_mode, ct_eff_map=ct_eff_map or {}, lang=lang, i18n=i18n, brand_primary=brand_primary, logo_path=logo_path)
    if vc_summary:
        add_value_chain_slide(prs, vc_summary, lang=lang, i18n=i18n, brand_primary=brand_primary, logo_path=logo_path)
    if material_flow_text:
        add_material_flow_narrative_slide(prs, material_flow_text, lang=lang, i18n=i18n, brand_primary=brand_primary, logo_path=logo_path)
    if finance:
        try:
            add_financial_slide(prs, finance, brand_primary=brand_primary, logo_path=logo_path)
        except Exception:
            pass
    if product_df is not None and champion is not None:
        try:
            add_product_selection_slide(prs, product_df, champion, brand_primary=brand_primary, logo_path=logo_path)
        except Exception:
            pass
    if savings:
        try:
            add_business_case_slide(prs, savings, brand_primary=brand_primary, logo_path=logo_path)
        except Exception:
            pass

    add_pqcdsm_slides(prs, observations_df, lang=lang, i18n=i18n, brand_primary=brand_primary, logo_path=logo_path)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _brand_header(slide, brand_primary, logo_path)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = tx.text_frame; tf.text = t_i18n("summary_top", lang, i18n or {})
    tf.paragraphs[0].font.size = Pt(28)
    y = 1.2
    for i, row in observations_df.head(8).iterrows():
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{row['step_name']} — {row['waste'].title()} | Score {row['score_0_5']:.1f} | RPN {row['rpn_pct']:.0f}% | {row.get('evidence','')}"
        p.font.size = Pt(14)
        y += 0.6

    for _, row in observations_df.iterrows():
        s = prs.slides.add_slide(prs.slide_layouts[5])
        _brand_header(s, brand_primary, logo_path)
        header = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        header.text_frame.text = f"{row['step_name']} — {row['waste'].title()}"
        header.text_frame.paragraphs[0].font.size = Pt(26)
        meta = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
        meta.text_frame.text = f"Score {row['score_0_5']:.1f} | RPN {row['rpn_pct']:.0f}% | Evidence: {row.get('evidence','')}"
        meta.text_frame.paragraphs[0].font.size = Pt(14)
        body = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6.5), Inches(3))
        body.text_frame.text = row['observation']
        body.text_frame.paragraphs[0].font.size = Pt(18)
        # Photos (up to 2) on right
        try:
            if photos:
                key = (str(row.get('step_id','')), str(row.get('waste','')).lower())
                files = photos.get(key, [])[:2]
                px = Inches(7.2); py = Inches(2.0); ph = Inches(1.9)
                for i, fp in enumerate(files):
                    if os.path.exists(fp):
                        s.shapes.add_picture(fp, px, py + Inches(i*2.1), height=ph)
        except Exception:
            pass

    prs.save(out_path)
    return out_path

def add_pqcdsm_slides(prs, observations_df, lang='en', i18n=None, brand_primary="#C00000", logo_path=None):
    theme_order = [("P","Production"),("Q","Quality"),("C","Cost"),("D","Delivery"),("S","Safety"),("M","Morale")]
    if "theme_code" not in observations_df.columns:
        def cat(w):
            w=str(w).lower()
            if w in ("overproduction","overprocessing","motion","transportation"): return "P"
            if w in ("defects",): return "Q"
            if w in ("inventory",): return "C"
            if w in ("waiting",): return "D"
            if w in ("safety",): return "S"
            if w in ("talent",): return "M"
            return "P"
        observations_df = observations_df.copy()
        observations_df["theme_code"] = observations_df["waste"].apply(cat)
    for code, name in theme_order:
        grp = observations_df[observations_df["theme_code"]==code]
        if grp.empty: 
            continue
        s = prs.slides.add_slide(prs.slide_layouts[5])
        _brand_header(s, brand_primary, logo_path)
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title.text_frame.text = f"{code} — {name}: {t_i18n('pqcdsm_obs', lang, i18n or {})}"
        title.text_frame.paragraphs[0].font.size = Pt(26)
        y = 1.2
        for idx, row in enumerate(grp.itertuples(index=False), start=1):
            tb = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
            tf = tb.text_frame; tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{code}-{idx} {row.step_name} — {row.waste.title()}"
            p.font.size = Pt(14); p.font.bold = True
            q = tf.add_paragraph()
            q.text = getattr(row, "observation", "")
            q.level = 1; q.font.size = Pt(12)
            y += 0.8
            if y > 6.5:
                y = 1.2
                s = prs.slides.add_slide(prs.slide_layouts[5])
                _brand_header(s, brand_primary, logo_path)

def add_current_state_map_slide(prs, steps, perstep_top2, spacing_mode="Effective CT", ct_eff_map=None, lang='en', i18n=None, brand_primary="#C00000", logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _brand_header(slide, brand_primary, logo_path)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title.text_frame.text = t_i18n("csm", lang, i18n or {})
    title.text_frame.paragraphs[0].font.size = Pt(28)
    info_y = Inches(0.9); mat_y = Inches(5.2)
    info_lane = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), info_y, Inches(9.2), Inches(0.5))
    info_lane.fill.solid(); info_lane.fill.fore_color.rgb = RGBColor(230,230,230)
    info_lane.text_frame.text = "Information Flow"
    mat_lane = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), mat_y, Inches(9.2), Inches(0.5))
    mat_lane.fill.solid(); mat_lane.fill.fore_color.rgb = RGBColor(230,230,230)
    mat_lane.text_frame.text = "Material Flow"
    start_x = Inches(0.5); box_w=Inches(2.6); box_h=Inches(2.4); box_y=Inches(2.0); gap = Inches(1.0)
    positions=[]; x=start_x
    for i,s in enumerate(steps):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_y, box_w, box_h)
        rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(198, 239, 206); rect.line.color.rgb = RGBColor(91,155,213)
        tf = rect.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = f"{s.id} – {s.name}"; p.font.size=Pt(14); p.font.bold=True
        for line in [f"CT: {int(s.ct_sec or 0)} s", f"WIP: {int(s.wip_units_in or 0)}", f"Defects: {getattr(s,'defect_pct',0):.1f}%", f"Mode: {getattr(s,'push_pull','')}"]:
            q = tf.add_paragraph(); q.text=line; q.level=1; q.font.size=Pt(11)
        positions.append((x, box_y)); x = x + box_w + gap
    for i in range(len(positions)-1):
        (x1,y1) = positions[i]; (x2,y2)=positions[i+1]
        conn = slide.shapes.add_connector(1, int(x1+box_w), int(y1+box_h/2), int(x2), int(y2+box_h/2))
        conn.line.width=Pt(2); conn.line.color.rgb = RGBColor(0,0,0)

def add_value_chain_slide(prs, vc_summary, lang='en', i18n=None, brand_primary="#C00000", logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _brand_header(slide, brand_primary, logo_path)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title.text_frame.text = t_i18n("vc", lang, i18n or {})
    title.text_frame.paragraphs[0].font.size = Pt(26)
    x = Inches(0.5); y = Inches(1.2); w = Inches(2.3); h = Inches(0.9); gap = Inches(0.4)
    for i, st in enumerate(vc_summary):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(230,242,255); box.line.color.rgb = RGBColor(91,155,213)
        tf = box.text_frame; tf.clear(); tf.text = st["stage_name"]; tf.paragraphs[0].font.size=Pt(12); tf.paragraphs[0].font.bold=True
        if i < len(vc_summary)-1:
            slide.shapes.add_connector(1, int(x+w), int(y+h/2), int(x+w+gap), int(y+h/2)).line.width=Pt(2)
        tb = slide.shapes.add_textbox(x, y+h+Inches(0.1), w, Inches(1.5))
        tf2 = tb.text_frame; tf2.clear(); tf2.text = "Top 3 Mudas"; tf2.paragraphs[0].font.size=Pt(11); tf2.paragraphs[0].font.bold=True
        for (m, sc) in st.get("top3", [])[:3]:
            p = tf2.add_paragraph(); p.text = f"• {m.title()} ({sc:.1f})"; p.font.size=Pt(10); p.level=1
        issues = st.get("issues", [])[:2]
        if issues:
            p = tf2.add_paragraph(); p.text = "Examples:"; p.font.size=Pt(10); p.level=1
            for iss in issues:
                q = tf2.add_paragraph(); q.text = f"– {iss}"; q.font.size=Pt(9); q.level=2
        x = x + w + gap

def export_observations_pdf(observations_df, out_path, brand_primary="#C00000", logo_path="assets/kafaa_logo.png"):
    c = canvas.Canvas(out_path, pagesize=landscape(A4))
    w, h = landscape(A4)

    # Watermark function
    def _watermark():
        try:
            from reportlab.lib.utils import ImageReader
            c.saveState()
            c.translate(w * 0.5, h * 0.35)
            c.rotate(25)
            try:
                c.setFillAlpha(0.08)  # available in reportlab 4.x
            except Exception:
                pass
            img = ImageReader(logo_path) if (logo_path and os.path.exists(logo_path)) else None
            if img:
                c.drawImage(img, -w*0.25, -h*0.15, width=w*0.5, height=h*0.3, preserveAspectRatio=True, mask='auto')
            c.restoreState()
        except Exception:
            pass

    _watermark()
    c.setFont("Helvetica-Bold", 20); c.drawString(2*cm, h-1.5*cm, "Automated VSM – Observations")
    c.setFont("Helvetica-Bold", 12); y = h-3.0*cm
    for _, row in observations_df.iterrows():
        c.drawString(1.5*cm, y, f"{row['step_name']} — {row['waste'].title()} (Score {row['score_0_5']:.1f} | RPN {row['rpn_pct']:.0f}% | {row.get('evidence','')})")
        y -= 0.7*cm; c.setFont("Helvetica", 11)
        for line in split_text(row['observation'], max_chars=140):
            c.drawString(2.0*cm, y, line); y -= 0.55*cm
            if y < 2.0*cm:
                c.showPage(); _watermark(); y = h-2.0*cm; c.setFont("Helvetica", 11)
        c.setFont("Helvetica-Bold", 12); y -= 0.3*cm
        if y < 3.0*cm:
            c.showPage(); _watermark(); y = h-3.0*cm; c.setFont("Helvetica-Bold", 12)
    c.showPage(); c.save(); return out_path

def split_text(text, max_chars=100):
    words = text.split(); out=[]; cur=""
    for w in words:
        if len(cur)+len(w)+1 <= max_chars:
            cur=(cur+" "+w).strip()
        else:
            out.append(cur); cur=w
    if cur:
        out.append(cur)
    return out


def add_financial_slide(prs, finance: dict, brand_primary="#C00000", logo_path=None):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _brand_header(slide, brand_primary, logo_path)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title.text_frame.text = "Financial Assessment — Targets"
    title.text_frame.paragraphs[0].font.size = Pt(26)

    # Left: summary metrics
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.2), Inches(4.5))
    tf = left.text_frame; tf.clear()
    def _p(text, lvl=0, size=14, bold=False):
        p = tf.add_paragraph() if len(tf.paragraphs)>0 else tf.paragraphs[0]
        p.text = text; p.level=lvl; p.font.size = Pt(size); p.font.bold = bold
        return p
    tf.text = f"Year: {finance.get('Year','-')}"
    tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.bold = True
    _p(f"Current Profit: {finance.get('current_profit_fmt','-')}")
    _p(f"Targeted Profit: {finance.get('Targeted Profit_fmt','-')}")
    _p(f"Profit Gap (actual): {finance.get('profit_gap_actual_fmt','-')}")
    _p(f"Req. cost reduction (to hit target at Sales Target): {finance.get('required_reduction_fmt','-')}", size=13)
    _p("Liquidity & inventory:", bold=True)
    _p(f"Quick Ratio (≈cash proxy): {finance.get('quick_ratio_str','-')}")
    _p(f"Inventory Days: {finance.get('inventory_days_str','-')}")
    _p(f"Inventory as % of Current Assets: {finance.get('inv_pct_ca_str','-')}")
    _p(f"Inventory reduction to reach Quick Ratio 1.0: {finance.get('inv_reduction_for_qr1_fmt','-')}", size=12)

    # Right: allocation
    right = slide.shapes.add_textbox(Inches(5.0), Inches(1.0), Inches(4.5), Inches(4.5))
    rf = right.text_frame; rf.clear()
    rf.text = "Suggested cost-reduction allocation"
    rf.paragraphs[0].font.size = Pt(14); rf.paragraphs[0].font.bold = True
    alloc = finance.get('allocation', {})
    for k in ['COGS','G&A','Financial Expenses','Depreciation']:
        if k in alloc:
            p = rf.add_paragraph(); p.level=1; p.text = f"{k}: {alloc[k]['amount_fmt']} ({alloc[k]['share']*100:.0f}%)"; p.font.size = Pt(12)
    if finance.get('notes'):
        p = rf.add_paragraph(); p.level=1; p.text = "Notes:"; p.font.size=Pt(12); p.font.bold=True
        for n in finance['notes'][:3]:
            q = rf.add_paragraph(); q.level=2; q.text = f"- {n}"; q.font.size=Pt(11)

    return slide


def add_product_selection_slide(prs, df, champion: dict, brand_primary="#C00000", logo_path=None):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _brand_header(slide, brand_primary, logo_path)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title.text_frame.text = "Product Selection — Champion SKU"
    title.text_frame.paragraphs[0].font.size = Pt(26)

    # Champion box
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.0))
    tf = box.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Champion: {champion.get('Product Name','-')}  |  Score: {champion.get('Total Score','-')}"
    p.font.size = Pt(18); p.font.bold = True

    # Top-5 table (name + score)
    top = df.sort_values("Total Score", ascending=False).head(5)
    x = Inches(0.5); y = Inches(2.2)
    for i, r in enumerate(top.itertuples(index=False), start=1):
        tb = slide.shapes.add_textbox(x, y, Inches(9), Inches(0.45))
        tf2 = tb.text_frame; tf2.clear()
        row = tf2.paragraphs[0]; row.font.size = Pt(12)
        row.text = f"{i}. {getattr(r,'Product Name','-')} — Score {getattr(r,'Total Score',0):.2f}"
        y += Inches(0.5)

    # Notes
    if "Notes" in champion:
        note = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.3))
        nt = note.text_frame; nt.text = f"Why chosen: {champion['Notes']}"
        nt.paragraphs[0].font.size = Pt(12)

    return slide


def export_charter_pdf(charter: dict, out_path, brand_primary="#C00000", logo_path="assets/kafaa_logo.png"):
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import cm
    from reportlab.lib.utils import ImageReader
    import os

    c = canvas.Canvas(out_path, pagesize=landscape(A4))
    w, h = landscape(A4)

    def _watermark():
        try:
            c.saveState(); c.translate(w*0.5, h*0.35); c.rotate(25)
            try:
                c.setFillAlpha(0.08)
            except Exception:
                pass
            if logo_path and os.path.exists(logo_path):
                c.drawImage(ImageReader(logo_path), -w*0.25, -h*0.15, width=w*0.5, height=h*0.3, preserveAspectRatio=True, mask='auto')
            c.restoreState()
        except Exception:
            pass

    def _section(title, y):
        c.setFillColorRGB(0,0,0); c.setFont("Helvetica-Bold", 14); c.drawString(1.5*cm, y, title)
        c.line(1.5*cm, y-0.15*cm, w-1.5*cm, y-0.15*cm)
        return y-0.6*cm

    _watermark()
    c.setFont("Helvetica-Bold", 20); c.drawString(2*cm, h-1.5*cm, "VSM Team Charter")
    c.setFont("Helvetica", 11)

    y = h-2.5*cm
    # Header fields
    hdr = [
        ("Value Stream Name", charter.get("vs_name","-")),
        ("Product", charter.get("product","-")),
        ("Starting Point", charter.get("start_point","-")),
        ("Ending Point", charter.get("end_point","-")),
        ("Workshop Location", charter.get("location","-")),
        ("Kick-off Date", charter.get("kickoff","-")),
    ]
    for i in range(0, len(hdr), 3):
        row = hdr[i:i+3]
        x = 1.5*cm; y -= 0.0
        for label, val in row:
            c.setFont("Helvetica-Bold", 11); c.drawString(x, y, f"{label}:")
            c.setFont("Helvetica", 11); c.drawString(x+3.8*cm, y, str(val))
            x += 8.2*cm
        y -= 0.7*cm

    # Roles
    y = _section("Roles", y)
    roles = [
        ("Executive Sponsor", charter.get("exec_sponsor","-")),
        ("Value Stream Owner", charter.get("owner","-")),
        ("Value Stream Champion", charter.get("champion_rep","-")),
        ("Workshop Facilitator", charter.get("facilitator","-")),
    ]
    for label, val in roles:
        c.setFont("Helvetica-Bold", 11); c.drawString(1.5*cm, y, f"{label}:")
        c.setFont("Helvetica", 11); c.drawString(6.2*cm, y, str(val)); y -= 0.55*cm

    # Objectives & KPIs
    y = _section("Objectives & Success Measures (KPIs)", y)
    txt = charter.get("objectives","").split("\n")
    for t in txt:
        if not t.strip(): continue
        c.drawString(1.8*cm, y, f"• {t.strip()}"); y -= 0.5*cm

    y = _section("Current State Issues & Business Needs", y)
    txt = charter.get("issues","").split("\n")
    for t in txt[:6]:
        if not t.strip(): continue
        c.drawString(1.8*cm, y, f"• {t.strip()}"); y -= 0.5*cm

    # Financial targets
    y = _section("Financial Targets", y)
    c.drawString(1.8*cm, y, f"Required Cost Reduction: {charter.get('required_reduction_fmt','-')}")
    y -= 0.5*cm
    c.drawString(1.8*cm, y, f"Quick Ratio (current): {charter.get('quick_ratio_str','-')}  |  Inventory Days: {charter.get('inventory_days_str','-')}")
    y -= 0.5*cm
    c.drawString(1.8*cm, y, f"Inventory reduction to QR=1.0: {charter.get('inv_reduction_for_qr1_fmt','-')}")

    # Team members (table)
    y = _section("Team Members", y)
    cols = ["Department","Name","Contact","Role","Signature"]
    x0 = 1.5*cm; colw = [4.0*cm, 5.0*cm, 5.0*cm, 4.0*cm, 4.0*cm]
    # header
    for i,(cname,wid) in enumerate(zip(cols,colw)):
        c.setFont("Helvetica-Bold", 10); c.drawString(x0, y, cname); x0 += wid
    y -= 0.35*cm
    c.line(1.5*cm, y, w-1.5*cm, y); y -= 0.2*cm
    # rows
    x0 = 1.5*cm
    for member in charter.get("team", [])[:7]:
        x = 1.5*cm
        for val,wid in zip([member.get("dept",""), member.get("name",""), member.get("contact",""), member.get("role",""), ""], colw):
            c.setFont("Helvetica", 10); c.drawString(x, y, str(val)[:40]); x += wid
        y -= 0.6*cm
        if y < 3.5*cm: break

    # Sign-off
    y = _section("Approvals", y)
    sign = [
        ("Executive Sponsor", charter.get("exec_sponsor","-")),
        ("Value Stream Owner", charter.get("owner","-")),
        ("Value Stream Champion", charter.get("champion_rep","-")),
        ("Date", charter.get("sign_date","-"))
    ]
    for label,val in sign:
        c.setFont("Helvetica-Bold", 11); c.drawString(1.5*cm, y, f"{label}:")
        c.setFont("Helvetica", 11); c.drawString(5.5*cm, y, str(val))
        c.line(5.4*cm, y-0.1*cm, 17.5*cm, y-0.1*cm); y -= 0.8*cm

    c.showPage(); c.save(); return out_path

def add_business_case_slide(prs, savings: dict, brand_primary="#C00000", logo_path=None):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _brand_header(slide, brand_primary, logo_path)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title.text_frame.text = "Business Case — Estimated Annual Impact"
    title.text_frame.paragraphs[0].font.size = Pt(26)

    bw = savings.get("by_waste", {})
    items = [(k, v) for k, v in bw.items() if v>0]
    items.sort(key=lambda x: x[1], reverse=True)
    y = Inches(1.0)
    total = savings.get("total", 0.0)
    box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.6))
    tf = box.text_frame; tf.text = f"Total Estimated Annual Benefit: {total:,.0f}"
    tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.bold = True
    y = Inches(1.8)
    for k,v in items[:6]:
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.45))
        tb.text_frame.text = f"• {k.title()}: {v:,.0f}"
        tb.text_frame.paragraphs[0].font.size = Pt(14)
        y += Inches(0.5)
    return slide
